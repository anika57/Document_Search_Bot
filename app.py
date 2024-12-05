from flask import Flask, request, jsonify, send_from_directory
import os
import google.generativeai as genai
from werkzeug.utils import secure_filename
from dotenv import load_dotenv
from flask_cors import CORS
import PyPDF2
import docx
import openpyxl
import pptx

# Initialize Flask app
app = Flask(__name__)

# Enable CORS for all routes
CORS(app)

# Load environment variables
load_dotenv()

# Load the API key from environment variables
api_key = os.getenv("GEMINI_API_KEY")
if not api_key:
    print("Error: Gemini API key is missing.")
    raise ValueError("Check Your API - Something Seems Wrong!!!")

# Configure the Gemini API
genai.configure(api_key=api_key)

generation_config = {
    "temperature": 1,
    "top_p": 0.95,
    "top_k": 64,
    "max_output_tokens": 500,
}

model = genai.GenerativeModel(
    model_name="gemini-1.5-flash",
    generation_config=generation_config,
)

chat_session = model.start_chat()

# Directory for storing uploaded files
UPLOAD_FOLDER = 'uploads'
if not os.path.exists(UPLOAD_FOLDER):
    os.makedirs(UPLOAD_FOLDER)

def extract_text_from_file(file_path, file_type):
    text = ""
    if file_type == 'pdf':
        with open(file_path, 'rb') as f:
            reader = PyPDF2.PdfReader(f)
            for page_num in range(len(reader.pages)):
                text += reader.pages[page_num].extract_text()
    elif file_type == 'docx':
        doc = docx.Document(file_path)
        for para in doc.paragraphs:
            text += para.text
    elif file_type == 'xlsx':
        wb = openpyxl.load_workbook(file_path)
        for sheet in wb:
            for row in sheet.iter_rows(values_only=True):
                text += ' '.join([str(cell) for cell in row if cell is not None]) + '\n'
    elif file_type == 'pptx':
        prs = pptx.Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text
    return text

def answer_query_from_text(text, query):
    prompt = (
        "You are a document assistant. Please answer the following query based on the provided document content:\n\n"
        f"Document Content:\n{text}\n\n"
        f"Query: {query}\n\n"
        "Answer:"
    )
    response = chat_session.send_message(prompt)
    return response.text

@app.route('/upload-and-query', methods=['POST'])
def upload_and_query():
    if 'file' not in request.files or 'query' not in request.form:
        return jsonify({"error": "File or query part is missing"}), 400
    
    file = request.files['file']
    query = request.form['query']
    
    if file.filename == '':
        return jsonify({"error": "No selected file"}), 400
    
    if file.mimetype not in ['application/pdf', 'application/vnd.openxmlformats-officedocument.wordprocessingml.document', 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet', 'application/vnd.openxmlformats-officedocument.presentationml.presentation']:
        return jsonify({"error": "Unsupported file type"}), 400
    
    if len(file.read()) > 2 * 1024 * 1024:
        return jsonify({"error": "File size exceeds 2MB"}), 400
    
    file.seek(0)
    filename = secure_filename(file.filename)
    file_path = os.path.join(UPLOAD_FOLDER, filename)
    
    try:
        file.save(file_path)
    except Exception as e:
        return jsonify({"error": f"Failed to save file: {str(e)}"}), 500

    try:
        file_type = file.mimetype.split('/')[-1]
        text = extract_text_from_file(file_path, file_type)
        answer = answer_query_from_text(text, query)
        return jsonify({"answer": answer.strip()}), 200

    except Exception as e:
        return jsonify({"error": f"Error processing file: {str(e)}"}), 500

@app.route('/list-files', methods=['GET'])
def list_files():
    files = os.listdir(UPLOAD_FOLDER)
    return jsonify({"files": files}), 200

@app.route('/delete-file', methods=['DELETE'])
def delete_file():
    file_name = request.args.get('file')
    if not file_name:
        return jsonify({"error": "File name not provided"}), 400
    
    file_path = os.path.join(UPLOAD_FOLDER, file_name)
    if not os.path.exists(file_path):
        return jsonify({"error": "File not found"}), 404
    
    try:
        os.remove(file_path)
        return jsonify({"message": f"File '{file_name}' deleted successfully"}), 200
    except Exception as e:
        return jsonify({"error": f"Failed to delete file: {str(e)}"}), 500

if __name__ == '__main__':
    app.run(debug=True)
